# -*- coding: utf-8 -*-
"""Cover + 8 deep-dive slides: conversational / agentic AI principles for SCB."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "workshop" / "SCB_Conversational_AI_Eight_Principles.pptx"

ACCENT = RGBColor(0x3F, 0x23, 0x6C)
BODY = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
CALL = RGBColor(0x00, 0x66, 0x99)


def add_textbox(slide, left, top, width, height, lines: list[tuple[str, int, bool, RGBColor | None]], font_default: int = 14):
    """lines: (text, size, bold, color or None for BODY)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color or BODY
        p.space_after = Pt(5)
        p.level = 0
        if text.startswith("    •") or text.startswith("•"):
            p.level = 1
            p.space_before = Pt(0)


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- Cover -----
    s = slide_blank(prs)
    add_textbox(
        s,
        Inches(0.6),
        Inches(0.55),
        Inches(12.0),
        Inches(1.1),
        [
            ("Conversational & agentic AI — eight CX principles", 32, True, ACCENT),
            ("Blueprint overview for SCB Bank", 18, False, MUTED),
        ],
        32,
    )
    cover_body = [
        ("Eight mutually reinforcing pillars — from first impression to continuous improvement.", 14, False, None),
        ("", 6, False, None),
        ("1. Personality & Voice — who speaks and how they earn credibility in stress moments.", 13, False, None),
        ("2. Intent Recognition & Context Awareness — messy language → correct intent without repetition loops.", 13, False, None),
        ("3. Conversational Flow Design — minimal friction paths, branching, and recovery for priority journeys.", 13, False, None),
        ("4. Personalization & Relevance — tailored help and nudges within explicit boundaries (incl. servicing vs sell).", 13, False, None),
        ("5. Transparency & Trust Design — disclosures, explainability, confirmations, and traceable actions.", 13, False, None),
        ("6. Feedback & Learning Loops — explicit + implicit signals feeding governed improvement.", 13, False, None),
        ("7. Human-in-the-Loop & Escalation — when humans enter and what context travels with the customer.", 13, False, None),
        ("8. Intuitiveness & Accessibility — natural language, languages, low cognitive load, inclusive multimodal UI.", 13, False, None),
        ("", 6, False, None),
        ("Next slides: one deep dive each + explicit design choices for SCB to decide.", 12, True, CALL),
    ]
    add_textbox(s, Inches(0.6), Inches(1.75), Inches(12.0), Inches(5.2), cover_body, 14)

    foot = s.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.45))
    foot.text_frame.text = "Grounded in banking CX practice, agentic-service patterns, and your servicing-first / deflection goals — tailor wording with SCB brand & legal."
    foot.text_frame.paragraphs[0].font.size = Pt(10)
    foot.text_frame.paragraphs[0].font.color.rgb = MUTED

    slides_data: list[dict] = [
        {
            "title": "1. Personality & Voice",
            "question": "Who is speaking? And why should I trust them?",
            "core": [
                "Defined tone (OCEAN)",
                "Consistency across journeys",
                "Aligned with brand values",
                "Emotionally calibrated (calm in stress, not overly cheerful in disputes)",
            ],
            "research": [
                "Industry shift: empathetic, state-aware service vs flat “happy bot” tone (McKinsey: empathetic CX with agentic AI).",
                "Trust is cumulative: voice mismatch after a failure is a top driver of abandonment in service chat.",
                "Design practice: persona = voice + boundaries (what the assistant will never joke about or claim).",
            ],
            "scb": [
                "Choose primary persona stance: expert-neutral vs warm-partner vs premium-concierge (one spine, not three).",
                "Define tone-by-state matrix for disputes, fraud-adjacent flows, and post-failed resolution.",
                "Set rules for humor, emojis, and colloquialisms by channel (chat vs voice vs push).",
                "Align Jarvis / assistant naming and disclosure timing with brand architecture (one character or family?).",
            ],
        },
        {
            "title": "2. Intent Recognition & Context Awareness",
            "question": "Does it understand me beyond my exact words?",
            "core": [
                "Accurate intent detection (even with messy, human phrasing)",
                "Memory of conversation context (no repetition loops)",
                "Ability to disambiguate gracefully",
            ],
            "research": [
                "Agentic systems amplify upstream errors: weak intent → wrong tool calls → higher remediation cost.",
                "Best practice: finite intent taxonomy for regulated journeys + robust “unknown” handling.",
                "Session memory contract: what persists for coherence vs what requires re-auth for security.",
            ],
            "scb": [
                "Prioritize intents for Journey 2: payment anomaly, loan servicing, card block (+ agreed extensions).",
                "Set confidence thresholds: when to clarify vs act vs escalate (document per intent).",
                "Policy for cross-session memory (e.g. last topic, preferences) vs fresh session for sensitive intents.",
                "Disambiguation UI: max options shown; fallbacks when ASR/NLP confidence is low (especially Thai/English mix if applicable).",
            ],
        },
        {
            "title": "3. Conversational Flow Design",
            "question": "Does the conversation move naturally and seamlessly?",
            "core": [
                "Clear pathways for key use cases (balance, disputes, loans, KYC)",
                "Minimal friction (no unnecessary steps)",
                "Smart branching logic (adapts based on responses)",
                "Recovery paths (when users go off-script)",
            ],
            "research": [
                "Service CX leaders design explicit exit, cancel, and “start over” primitives — reduces silent drop-off.",
                "Agentic flows need visible milestones: users tolerate multi-step automation when progress is legible.",
                "Friction budget: cap clarifying turns before offering human or alternate path (industry service-design norm).",
            ],
            "scb": [
                "One canonical flow diagram per priority journey; align with call-center deflection targets.",
                "Define max steps / max clarifications before mandatory escalation or callback offer.",
                "Standard recovery patterns: wrong account, wrong amount, user silence/timeout, authenticated vs guest paths.",
                "Where conversational handoff meets app UI (deeplinks, forms): who owns the next step?",
            ],
        },
        {
            "title": "4. Personalization & Relevance",
            "question": "Does it feel like it knows me?",
            "core": [
                "Use of customer data",
                "Contextual nudges",
                "Dynamic responses",
                "Respectful boundaries",
            ],
            "research": [
                "Deloitte-type finding pattern: personalization lifts engagement but trust for advice remains fragile — justify relevance.",
                "BCG / market direction: agent-led engagement must not feel like surveillance; frequency caps reduce annoyance.",
                "Your workshop context: research tension on offers after servicing — commercial posture must be explicit policy.",
            ],
            "scb": [
                "Write NBA / cross-sell rules: suppressed / deferred / allowed windows relative to service outcome and sentiment.",
                "Data-use ladder: what is used for copy only vs ranking options vs pre-filling amounts (with confirmation).",
                "Caps on nudges per day/week; rules for stacking (no promo pile-on after failed resolution).",
                "Segment rules vs true 1:1; when “why am I seeing this?” copy is mandatory.",
            ],
        },
        {
            "title": "5. Transparency & Trust Design",
            "question": "Can I see what’s happening?",
            "core": [
                "Clear disclosures (“I’m a virtual assistant”)",
                "Explainability (“This charge is from…”)",
                "Confirmation before critical actions",
                "Auditability (traceable interactions)",
            ],
            "research": [
                "Regulated and high-trust sectors converge on: capability limits stated plainly, tiered confirmations, reconstructable logs.",
                "Agentic AI: users need plan transparency at checkpoints — what will happen before it happens.",
                "ING-style pattern (public case literature): retrieval + ranking + guardrails as part of customer-visible quality.",
            ],
            "scb": [
                "Disclosure stack: identity, non-human, non-legal-advice where relevant, data sources for explanations.",
                "Risk-tiered confirmation map: view vs transfer vs dispute submission vs hardship restructuring request.",
                "What the assistant may never claim (e.g. final credit decision, legal outcome) — hard refusal templates.",
                "Audit record: minimum fields for complaints handling; alignment with internal ops and regulator expectations.",
            ],
        },
        {
            "title": "6. Feedback & Learning Loops",
            "question": "Does the system get better over time?",
            "core": [
                "Explicit feedback (ratings, thumbs up/down)",
                "Implicit signals (drop-offs, retries, escalations)",
                "Continuous training pipelines",
                "Closed-loop improvement (feedback → action → refinement)",
            ],
            "research": [
                "Leading programs pair qualitative review with metrics — containment without quality is a false win.",
                "Governed change: which intents require human approval before prompt/model updates (bank norm).",
                "Closed loop: owner, hypothesis, release, and post-release monitoring per intent cluster.",
            ],
            "scb": [
                "Metric suite: resolution quality, recontact, escalation reason codes, CSAT by intent — not chat volume alone.",
                "Where to ask for explicit feedback (milestone-based) vs lightweight thumb to avoid fatigue.",
                "RACI for prompt / flow / model updates: risk, legal, brand sign-off by journey tier.",
                "How ACX app research and conversational telemetry feed one backlog (avoid two truths).",
            ],
        },
        {
            "title": "7. Human-in-the-Loop & Escalation",
            "question": "Are handoffs to humans seamless and context-complete?",
            "core": [
                "Seamless handoff to human agents",
                "Context transfer (no “repeat everything”)",
                "Smart escalation triggers (emotion, complexity, risk)",
                "Blended experiences (bot assists human behind the scenes)",
            ],
            "research": [
                "BCG / service transformation theme: orchestration across AI + human — handoff is a designed product, not a failure state.",
                "Customer effort score spikes hardest on “repeat myself” — warm transfer packet is table stakes.",
                "Triggers: confidence, sentiment, regulatory path, fraud signals, explicit customer ask.",
            ],
            "scb": [
                "Escalation trigger list with owners: e.g. fraud, legal threat, third-party payment, hardship escalation.",
                "Warm handoff data packet schema (customer-visible summary + agent CRM fields).",
                "When bot must stop talking vs co-pilot for agent (compliance-checked scripts only).",
                "Alignment with call-center migration / deflection program — no incentive to trap users in bot.",
            ],
        },
        {
            "title": "8. Intuitiveness & Accessibility",
            "question": "Can anyone use this without thinking twice?",
            "core": [
                "Natural language inputs (no command training required)",
                "Multilingual support",
                "Low cognitive load (clear prompts, no jargon)",
                "Accessible design (voice, readability, inclusivity)",
                "Visual attachments (rich media)",
            ],
            "research": [
                "WCAG-aligned patterns: not color-only status, focus order, labels on chips, text alternatives for critical visuals.",
                "Multilingual: plan for unsupported locales (clear fallback, not silent degradation).",
                "Rich media: one primary action per card; summarize charts for voice-only or low-vision paths.",
                "Design critique you cited: density control — hierarchy and progressive disclosure beat “shorter text” alone.",
            ],
            "scb": [
                "Supported language list + behavior for code-switching / mixed utterances in market.",
                "Reading level target and banned jargon list for money and legal concepts.",
                "Voice vs chat parity: what must be speakable vs screen-only; when to force secure channel.",
                "Rich pattern library (cards, lists, charts) with accessibility acceptance criteria before ship.",
            ],
        },
    ]

    for block in slides_data:
        s = slide_blank(prs)
        y = Inches(0.45)
        add_textbox(
            s,
            Inches(0.55),
            y,
            Inches(12.0),
            Inches(0.55),
            [(block["title"], 26, True, ACCENT)],
        )
        y2 = Inches(1.05)
        add_textbox(
            s,
            Inches(0.55),
            y2,
            Inches(12.0),
            Inches(0.45),
            [(block["question"], 16, True, CALL)],
        )
        y3 = Inches(1.5)
        lines = [("Core principles", 12, True, MUTED), ("", 8, False, None)]
        for b in block["core"]:
            lines.append((f"    • {b}", 13, False, None))
        lines += [("", 8, False, None), ("Research & design practice", 12, True, MUTED)]
        for b in block["research"]:
            lines.append((f"    • {b}", 12, False, None))
        lines += [("", 8, False, None), ("Design choices for SCB Bank to decide", 12, True, ACCENT)]
        for b in block["scb"]:
            lines.append((f"    • {b}", 12, False, BODY))
        add_textbox(s, Inches(0.55), y3, Inches(12.0), Inches(5.7), lines, 13)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
