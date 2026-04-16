# -*- coding: utf-8 -*-
"""Generate SCB Conversational Banking CX Blueprint workshop deck (python-pptx)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "workshop" / "SCB_Conversational_CX_Blueprint_Workshop.pptx"

ACCENT = RGBColor(0x3F, 0x23, 0x6C)  # purple bank–adjacent
BODY = RGBColor(0x2B, 0x2B, 0x2B)
MUTED = RGBColor(0x55, 0x55, 0x55)


def _set_title(tf, text: str, size: int = 28):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.LEFT


def _add_body(slide, left, top, width, height, lines: list[str], size: int = 14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = BODY if not line.startswith("→") else MUTED
        p.space_after = Pt(6)
        p.level = 0
        if line.startswith("  •") or line.startswith("•"):
            p.level = 1


def add_slide(prs, title: str, body_lines: list[str], subtitle: str | None = None):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9.0), Inches(0.9))
    _set_title(title_box.text_frame, title, 26)
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(9.0), Inches(0.45))
        st.text_frame.text = subtitle
        st.text_frame.paragraphs[0].font.size = Pt(12)
        st.text_frame.paragraphs[0].font.color.rgb = MUTED
    _add_body(slide, Inches(0.55), Inches(1.45 if subtitle else 1.2), Inches(9.0), Inches(5.5), body_lines, 15)
    return slide


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(11.5), Inches(1.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Conversational banking CX blueprint"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = "Workshop — align on design principles & choices"
    p2.font.size = Pt(22)
    p2.font.color.rgb = BODY
    p3 = tf.add_paragraph()
    p3.text = "SCB · 2.5 hours · interactive"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED
    foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(11.5), Inches(0.6))
    foot.text_frame.text = "Internal working session — tailor client-facing wording with SCB brand & legal."
    foot.text_frame.paragraphs[0].font.size = Pt(11)
    foot.text_frame.paragraphs[0].font.color.rgb = MUTED

    add_slide(
        prs,
        "Why this workshop (product + design lens)",
        [
            "Lock a bounded CX blueprint for conversational banking — not full app IA.",
            "Bring external signal (industry + UX research) into deliberate choices SCB still needs to make.",
            "Reduce late rework: agree principles first; reflect in journeys & UI in the following weeks.",
            "Illustrate with prioritized servicing journeys (e.g. payment anomaly, loan servicing, card block) + near-term launches.",
        ],
        "Grounded in your prior alignment: hypothesis-led convergence, not open ideation on use-case selection.",
    )

    add_slide(
        prs,
        "Who is in the room (from working session context)",
        [
            "~6–8 participants · conversational AI owners (business + tech).",
            "Digital / ACX — mobile app evolution & user testing readouts.",
            "Design — component language, density, tone in UI.",
            "Servicing / contact center stakeholders where deflection & loan flows matter.",
            "McKinsey — perspective, facilitation, synthesis to decisions.",
        ],
        "Confirm names & single decision owner per stream before the session.",
    )

    add_slide(
        prs,
        "Outcomes by end of 2.5 hours",
        [
            "1) Shared language: 8 blueprint pillars + anchor question each.",
            "2) Alignment: where SCB agrees today vs “decide by [date]” vs needs evidence.",
            "3) Three workshop artifacts: priority matrix, escalation/NBA guardrails sketch, action log with owners.",
            "4) Input to UI library / prototypes: tone, confirmation patterns, handoff data packet.",
        ],
    )

    add_slide(
        prs,
        "Agenda — 2.5 hours (150 min) — interactive",
        [
            "00:00–00:15  Opening · objectives · parking lot · norms (cameras, decisions, Chatham House if needed).",
            "00:15–00:35  Industry pulse — 3 themes + 2-min neighbor chat after each slide.",
            "00:35–00:55  SCB playback — what is decided / tested / in-flight (pre-brief or live 15 min).",
            "00:55–01:25  Gallery + dot vote — maturity vs importance for each pillar (digital board or sticky dots).",
            "01:25–01:35  Break.",
            "01:35–02:15  Breakouts (3 pods) → rotate “red team” on one other pod’s output (10 min).",
            "02:15–02:40  Convergence — decision sheet, owners, evidence gaps.",
            "02:40–02:50  Close — comms, pre-reads, follow-up with research artifacts.",
        ],
        "Facilitator keeps a visible timer; parking lot for scope creep on whole-app redesign.",
    )

    add_slide(
        prs,
        "External signal — McKinsey (recent themes)",
        [
            "Agentic AI shifts service from reactive answers to orchestrated resolution — CX must design for handoffs, verification, and audit.",
            "Empathetic, context-aware experiences: calibrate tone and action to customer state — not one flat “happy bot.”",
            "Mobile + conversational stack: world-class apps pair clarity, speed, and trust moments at high-stakes steps.",
            "→ Read: Beyond the bot — empathetic CX with agentic AI (McKinsey Operations).",
            "→ Read: Agentic AI in banking — frontline readiness (McKinsey Financial Services).",
        ],
        "https://www.mckinsey.com/capabilities/operations/our-insights/beyond-the-bot-building-empathetic-customer-experiences-with-agentic-ai",
    )

    add_slide(
        prs,
        "External signal — BCG & Deloitte (market + customer reality)",
        [
            "BCG: retail banking value increasingly tied to agent-led service & sales — design for measurable conversion and responsible automation.",
            "BCG: customer service transformation = orchestration across human + AI, not chat-only cost takeout.",
            "Deloitte (2025 US survey): chatbots are common but trust for advice remains low — transparency, verification, and proven outcomes matter.",
            "Design implication: earn trust progressively; do not oversell assistant capability.",
            "→ BCG: Branches to Bots; New frontier in customer service transformation.",
            "→ Deloitte Insights: AI banking chatbots — from frustration to delight.",
        ],
        "Use as stimulus, not prescription — SCB choices still drive the blueprint.",
    )

    add_slide(
        prs,
        "Interactive — Dot vote + maturity grid (35 min block)",
        [
            "Axis Y: Importance to SCB conversational strategy (next 6 months).",
            "Axis X: Current clarity (immature → mature).",
            "Each participant: 5 votes total across the 8 pillars (dots or digital poll).",
            "Facilitator: circle top-right (important + immature) = workshop depth; bottom-left = defer.",
            "Output photo / export — becomes appendix A in workshop summary.",
        ],
    )

    add_slide(
        prs,
        "Breakouts (40 min) — 3 pods + red-team rotation",
        [
            "Pod A — Trust & safety: disclosures, confirmations, explainability, audit, model limits.",
            "Pod B — Flow & intelligence: intent, context, branching, recovery, servicing-first vs NBA timing.",
            "Pod C — Persona & inclusion: voice, multilingual, accessibility, rich media without cognitive overload.",
            "Each pod: 1 poster with — (1) 3 non-negotiables (2) 2 tradeoffs (3) 1 experiment to run with customers.",
            "Red team (10 min): another pod stress-tests using a “frustrated post-dispute customer” scenario card.",
        ],
    )

    pillars = [
        (
            "1 · Persona, voice & emotional calibration",
            "Anchor: Who is speaking — and how should they sound when the customer is calm vs frustrated vs in dispute?",
            [
                "Choices: persona definition; tone matrix by state; channel fit (chat / voice / push).",
                "Align with brand; avoid toxic positivity after failed transactions.",
                "Cross-cutting: commercial posture — when NOT to pitch (e.g. post-servicing stress).",
            ],
        ),
        (
            "2 · Intent, context & grounding",
            "Anchor: How do we know what they want when language is messy — without overstepping privacy?",
            [
                "Robust intent + entity extraction; graceful disambiguation prompts.",
                "Session context vs durable personalization (separate policies).",
                "Ground answers in permitted data sources; confidence thresholds for action.",
            ],
        ),
        (
            "3 · Journey & conversational flow architecture",
            "Anchor: What is the happy path, recovery path, and clean exit for each priority journey?",
            [
                "Minimal friction; smart branching; explicit back / cancel / start-over.",
                "Maps to servicing journeys: balance & anomalies, loans, card controls, KYC where applicable.",
            ],
        ),
        (
            "4 · Personalization & relevance (with boundaries)",
            "Anchor: What should feel personalized — and what requires consent or explanation?",
            [
                "Contextual nudges; dynamic content; frequency & relevance caps.",
                "“Why am I seeing this?” for any proactive offer or pre-filled suggestion.",
            ],
        ),
        (
            "5 · Transparency, verification & trust",
            "Anchor: What must we always disclose, explain, and confirm before money or data moves?",
            [
                "Assistant identity; limits of capability; plain-language reasons for charges.",
                "High-stakes confirmations; immutable audit trail for regulated actions.",
            ],
        ),
        (
            "6 · Feedback & governed learning loops",
            "Anchor: How do we improve weekly without bypassing risk review?",
            [
                "Explicit (thumbs, CSAT) + implicit (retries, escalations, silent exits).",
                "Human-in-the-loop for training changes on regulated intents.",
            ],
        ),
        (
            "7 · Human handoff & escalation",
            "Anchor: When do we transfer — and what context packet prevents “repeat everything”?",
            [
                "Triggers: sentiment, complexity, fraud risk, regulatory paths.",
                "Warm transfer: structured summary + customer-visible status.",
            ],
        ),
        (
            "8 · Inclusion, multimodal UX & accessibility",
            "Anchor: How do we keep cognitive load low for every segment and modality?",
            [
                "Natural language first; multilingual roadmap tied to segments.",
                "Readable hierarchy; voice parity; charts/cards that screen readers can parse.",
            ],
        ),
    ]

    add_slide(
        prs,
        "The eight blueprint pillars (Akshata backbone + industry sharpening)",
        [
            "Eight deliberate pillars — each becomes one row in the CX blueprint decision sheet.",
            "Principles = non-negotiables; choices = SCB-specific stance (A vs B) with rationale.",
            "Next slides: one pillar per slide with anchor question + workshop prompts.",
        ],
    )

    for title, anchor, bullets in pillars:
        lines = [anchor, ""] + ["• " + b for b in bullets] + ["", "→ Live prompt: 90 sec — what is already decided at SCB? What evidence exists?"]
        add_slide(prs, title, lines)

    add_slide(
        prs,
        "Convergence template (project onto screen)",
        [
            "For each pillar: Agree / Agree with edits / Needs data — pick one.",
            "Capture: Decision statement (one sentence) · Owner · Evidence or pilot needed · Link to journey demo.",
            "Explicit “non-goals” for this program phase (reduces scope creep).",
        ],
    )

    add_slide(
        prs,
        "Materials & logistics checklist",
        [
            "Pre-read: sanitized conversational UI library + any user-research summaries.",
            "Room: three breakout spaces or virtual rooms; timer; shared Miro / FigJam / Slides.",
            "Roles: facilitator + timekeeper + note-taker (rotates each block).",
            "Post-workshop: 48-hour summary PDF + updated blueprint v0.9.",
        ],
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
