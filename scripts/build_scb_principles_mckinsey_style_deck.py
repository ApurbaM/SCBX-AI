# -*- coding: utf-8 -*-
"""
SCB conversational / agentic AI principles deck — visual system aligned to
SCB_CXO_Board_Dashboard.html (CEO / CXO discussion cockpit): dark canvas, gold
labels, purple accents, panel-style content blocks.

Output: workshop/SCB_Conversational_AI_Principles_McKinsey_Style.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "workshop" / "SCB_Conversational_AI_Principles_McKinsey_Style.pptx"
OUT_AGENDA_ONLY = Path(__file__).resolve().parent.parent / "workshop" / "SCB_CX_Workshop_Agenda_2_5hr_OnePage.pptx"
OUT_PRINCIPLES_Q_ONLY = Path(__file__).resolve().parent.parent / "workshop" / "SCB_Eight_Principles_Guiding_Questions.pptx"

# SCB_CXO_Board_Dashboard.html :root tokens (approximate for RGB slides)
BG_DEEP = RGBColor(0x05, 0x04, 0x0F)
BG_PANEL = RGBColor(0x12, 0x0C, 0x22)
TEXT = RGBColor(0xF4, 0xF0, 0xFF)
TEXT_MUTED = RGBColor(0xA8, 0xA0, 0xC0)
GOLD = RGBColor(0xE8, 0xD5, 0xA0)
PURPLE = RGBColor(0x7B, 0x5E, 0xA7)
STROKE = RGBColor(0x55, 0x4A, 0x6E)
FONT = "Arial"


def _set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _gold_rule(slide, top, width=Inches(12.2), left=Inches(0.55)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def _panel(slide, left, top, width, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_PANEL
    shp.line.color.rgb = STROKE
    shp.line.width = Pt(1)
    try:
        shp.adjustments[0] = 0.08  # corner radius (rounded rect)
    except (AttributeError, IndexError, ValueError):
        pass
    return shp


def _apply_font_to_paragraph(p):
    p.font.name = FONT
    if p.runs:
        for r in p.runs:
            r.font.name = FONT


def _fill_text_frame(
    tf,
    paragraphs: list[tuple[str, int, bool, RGBColor, int]],
    default_color=TEXT,
):
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    first = True
    for text, size, bold, color, level in paragraphs:
        if not str(text).strip() and first:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color or default_color
        p.space_after = Pt(5 if level > 0 else 7)
        p.line_spacing = 1.18
        p.alignment = PP_ALIGN.LEFT
        _apply_font_to_paragraph(p)


def slide_deep_dive(
    prs,
    index: int,
    lede: str,
    question: str,
    insight: str,
    left_paras: list[tuple[str, int, bool, RGBColor | None, int]],
    right_paras: list[tuple[str, int, bool, RGBColor | None, int]],
):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, BG_DEEP)

    # Top gold accent (CXO header pattern)
    _gold_rule(s, Inches(0.38))

    # Eyebrow — principle index
    ey = s.shapes.add_textbox(Inches(0.6), Inches(0.52), Inches(12.0), Inches(0.35))
    _fill_text_frame(
        ey.text_frame,
        [
            (
                f"CONVERSATIONAL CX PRINCIPLE  ·  {index} OF 8",
                10,
                True,
                GOLD,
                0,
            ),
        ],
        TEXT,
    )
    ey.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # LEDE — principle name (primary hook)
    le = s.shapes.add_textbox(Inches(0.6), Inches(0.82), Inches(12.0), Inches(0.95))
    _fill_text_frame(le.text_frame, [(lede, 30, True, TEXT, 0)], TEXT)
    le.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Guiding question
    qq = s.shapes.add_textbox(Inches(0.6), Inches(1.72), Inches(12.0), Inches(0.42))
    _fill_text_frame(
        qq.text_frame,
        [(question, 15, False, PURPLE, 0)],
        TEXT,
    )

    # Insight line (former “action title”) — supporting lede
    ins = s.shapes.add_textbox(Inches(0.6), Inches(2.12), Inches(12.0), Inches(0.38))
    _fill_text_frame(
        ins.text_frame,
        [(insight, 12, True, TEXT_MUTED, 0)],
        TEXT,
    )

    y0 = Inches(2.52)
    col_h = Inches(4.55)
    w_l = Inches(6.05)
    w_r = Inches(5.95)
    gap = Inches(0.18)
    lx = Inches(0.55)
    rx = lx + w_l + gap
    pad = Inches(0.22)

    _panel(s, lx, y0, w_l, col_h)
    _panel(s, rx, y0, w_r, col_h)

    # Column titles (gold caps — CEO doc .persona-block-label pattern)
    lh = s.shapes.add_textbox(lx + pad, y0 + Inches(0.14), w_l - 2 * pad, Inches(0.36))
    _fill_text_frame(
        lh.text_frame,
        [("WHAT THIS MEANS — NUANCE, EXAMPLES, SCB CHOICES", 9, True, GOLD, 0)],
        TEXT,
    )
    rh = s.shapes.add_textbox(rx + pad, y0 + Inches(0.14), w_r - 2 * pad, Inches(0.36))
    _fill_text_frame(
        rh.text_frame,
        [("WHERE LEADING RETAIL APPS EXEMPLIFY (ILLUSTRATIVE)", 9, True, GOLD, 0)],
        TEXT,
    )

    bl = s.shapes.add_textbox(lx + pad, y0 + Inches(0.48), w_l - 2 * pad, col_h - Inches(0.62))
    _fill_text_frame(bl.text_frame, left_paras, TEXT)

    br = s.shapes.add_textbox(rx + pad, y0 + Inches(0.48), w_r - 2 * pad, col_h - Inches(0.62))
    _fill_text_frame(br.text_frame, right_paras, TEXT)

    foot = s.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.32))
    _fill_text_frame(
        foot.text_frame,
        [
            (
                "Illustrative competitive references — verify current product behavior before external use. Visual system mirrors SCB CXO board cockpit for CEO-style discussion.",
                8,
                False,
                TEXT_MUTED,
                0,
            )
        ],
        TEXT_MUTED,
    )


def slide_agenda_one_page(prs: Presentation):
    """Single-slide 2.5h interactive workshop agenda (CXO visual system)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, BG_DEEP)
    _gold_rule(s, Inches(0.38))

    k = s.shapes.add_textbox(Inches(0.6), Inches(0.48), Inches(12.0), Inches(0.35))
    _fill_text_frame(
        k.text_frame,
        [("SCB  ·  CXO / CEO DISCUSSION SERIES  ·  WORKSHOP RUNNER", 10, True, GOLD, 0)],
        TEXT,
    )

    title = s.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(12.0), Inches(0.55))
    _fill_text_frame(
        title.text_frame,
        [("Workshop agenda — 2.5 hours (150 min) · Interactive", 26, True, TEXT, 0)],
        TEXT,
    )

    sub = s.shapes.add_textbox(Inches(0.6), Inches(1.28), Inches(12.0), Inches(0.38))
    _fill_text_frame(
        sub.text_frame,
        [
            (
                "Objective: align on eight conversational / agentic CX principles — decisions, evidence gaps, and 48-hour actions.",
                12,
                False,
                PURPLE,
                0,
            ),
        ],
        TEXT,
    )

    _panel(s, Inches(0.55), Inches(1.72), Inches(12.2), Inches(5.15))
    body = s.shapes.add_textbox(Inches(0.75), Inches(1.88), Inches(11.85), Inches(4.85))
    rows: list[tuple[str, int, bool, RGBColor, int]] = [
        ("0:00–0:10  ·  Opening, objectives, norms", 11, True, TEXT, 0),
        ("    Interactive — one-word / sticky wall: “What makes today a win?”  ·  Parking lot + decision norms.", 10, False, TEXT_MUTED, 1),
        ("", 3, False, TEXT, 0),
        ("0:10–0:25  ·  Pulse on the eight principles", 11, True, TEXT, 0),
        ("    Interactive — draw a principle #1–8; 30s “in customer words”; neighbor +1 only  ·  Map strong vs needs-work.", 10, False, TEXT_MUTED, 1),
        ("", 3, False, TEXT, 0),
        ("0:25–0:45  ·  SCB reality drop (ACX / research / servicing)", 11, True, TEXT, 0),
        ("    Interactive — tag each claim E (evidence) or B (belief)  ·  Build follow-up evidence list.", 10, False, TEXT_MUTED, 1),
        ("", 3, False, TEXT, 0),
        ("0:45–1:15  ·  Dot vote + heatmap", 11, True, TEXT, 0),
        ("    Interactive — 5 dots / person on Importance × Clarity  ·  Optional 5-min pre-mortem on top cluster.", 10, False, TEXT_MUTED, 1),
        ("", 3, False, TEXT, 0),
        ("1:15–1:25  ·  Break", 11, True, TEXT, 0),
        ("", 3, False, TEXT, 0),
        ("1:25–2:05  ·  Breakouts + red team", 11, True, TEXT, 0),
        ("    Interactive — 3 pods → poster (3 non-negotiables, 2 tradeoffs, 1 test)  ·  Rotate posters; red-team with “angry post-dispute + loan pitch” card.", 10, False, TEXT_MUTED, 1),
        ("", 3, False, TEXT, 0),
        ("2:05–2:35  ·  Gallery + convergence", 11, True, TEXT, 0),
        ("    Interactive — 90s / pod gallery  ·  Decision sheet: Agree / Agree w/ edits / Needs data + owner + date  ·  Fist-of-five on blueprint completeness.", 10, False, TEXT_MUTED, 1),
        ("", 3, False, TEXT, 0),
        ("2:35–2:50  ·  Commitments & close", 11, True, TEXT, 0),
        ("    Interactive — each leader: one 48-hour commitment (verb + object + date)  ·  Confirm summary + artifacts within 48h.", 10, False, TEXT_MUTED, 1),
        ("", 4, False, TEXT, 0),
        ("Materials: visible timer · shared board (Miro/FigJam/Slides) · scenario card · pod half-sheet template · optional Mentimeter for open polls.", 9, True, PURPLE, 0),
    ]
    _fill_text_frame(body.text_frame, rows, TEXT)

    foot = s.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.38))
    _fill_text_frame(
        foot.text_frame,
        [
            (
                "~6–10 participants typical · Facilitator protects breakout + convergence time · Adjust playback length if research walkthrough runs over.",
                8,
                False,
                TEXT_MUTED,
                0,
            ),
        ],
        TEXT_MUTED,
    )


def slide_principles_and_questions(prs: Presentation, principles_list: list[dict]):
    """One slide: all eight principle ledes + guiding questions (matches deep-dive slides)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, BG_DEEP)
    _gold_rule(s, Inches(0.38))

    k = s.shapes.add_textbox(Inches(0.6), Inches(0.48), Inches(12.0), Inches(0.32))
    _fill_text_frame(
        k.text_frame,
        [("SCB  ·  CXO / CEO DISCUSSION SERIES  ·  PRINCIPLE MAP", 10, True, GOLD, 0)],
        TEXT,
    )

    tit = s.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(12.0), Inches(0.52))
    _fill_text_frame(
        tit.text_frame,
        [("Eight principles — guiding questions (deep dive map)", 24, True, TEXT, 0)],
        TEXT,
    )

    _panel(s, Inches(0.55), Inches(1.38), Inches(12.2), Inches(5.75))
    half_w = Inches(5.75)
    pad = Inches(0.28)
    lx = Inches(0.55) + pad
    rx = Inches(0.55) + half_w + Inches(0.35) + pad
    y_top = Inches(1.52)
    col_h = Inches(5.45)

    def col_paras(start_idx: int, end_idx: int) -> list[tuple[str, int, bool, RGBColor, int]]:
        paras: list[tuple[str, int, bool, RGBColor, int]] = []
        for idx in range(start_idx, end_idx):
            p = principles_list[idx]
            n = idx + 1
            paras.append((f"{n} · {p['lede']}", 11, True, TEXT, 0))
            paras.append((f"    {p['q']}", 10, False, PURPLE, 1))
            paras.append(("", 4, False, TEXT, 0))
        return paras

    left = s.shapes.add_textbox(lx, y_top, half_w - pad, col_h)
    _fill_text_frame(left.text_frame, col_paras(0, 4), TEXT)

    right = s.shapes.add_textbox(rx, y_top, half_w - pad, col_h)
    _fill_text_frame(right.text_frame, col_paras(4, 8), TEXT)

    foot = s.shapes.add_textbox(Inches(0.55), Inches(7.18), Inches(12.2), Inches(0.28))
    _fill_text_frame(
        foot.text_frame,
        [
            (
                "Each following deep-dive slide uses the same lede + guiding question, plus insight line and two-column detail.",
                8,
                False,
                TEXT_MUTED,
                0,
            ),
        ],
        TEXT_MUTED,
    )


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- Cover (CXO cockpit style) -----
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s0, BG_DEEP)
    _gold_rule(s0, Inches(0.42))

    kicker = s0.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12.0), Inches(0.4))
    _fill_text_frame(
        kicker.text_frame,
        [("SCB  ·  CXO / CEO DISCUSSION SERIES", 11, True, GOLD, 0)],
        TEXT,
    )

    t = s0.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(11.8), Inches(1.0))
    _fill_text_frame(
        t.text_frame,
        [
            ("Conversational & agentic AI", 34, True, TEXT, 0),
            ("Eight principles — deep dive", 20, False, TEXT_MUTED, 0),
        ],
        TEXT,
    )

    sub = s0.shapes.add_textbox(Inches(0.6), Inches(2.15), Inches(11.5), Inches(0.55))
    _fill_text_frame(
        sub.text_frame,
        [
            (
                "Retail banking CX blueprint  ·  Nuances, examples, best-in-class references, and explicit SCB decisions",
                13,
                False,
                PURPLE,
                0,
            ),
        ],
        TEXT,
    )

    _panel(s0, Inches(0.55), Inches(2.85), Inches(12.2), Inches(3.95))
    body = s0.shapes.add_textbox(Inches(0.82), Inches(3.05), Inches(11.65), Inches(3.55))
    cover_lines = [
        ("The eight principles (each is the lede of its own slide):", 10, True, GOLD, 0),
        ("", 4, False, TEXT, 0),
        ("1  Personality & Voice", 12, True, TEXT, 1),
        ("2  Intent Recognition & Context Awareness", 12, True, TEXT, 1),
        ("3  Conversational Flow Design", 12, True, TEXT, 1),
        ("4  Personalization & Relevance", 12, True, TEXT, 1),
        ("5  Transparency & Trust Design", 12, True, TEXT, 1),
        ("6  Feedback & Learning Loops", 12, True, TEXT, 1),
        ("7  Human-in-the-Loop & Escalation", 12, True, TEXT, 1),
        ("8  Intuitiveness & Accessibility", 12, True, TEXT, 1),
        ("", 6, False, TEXT, 0),
        (
            "Next: framing → eight principles + guiding questions (map slide) → eight deep dives → sources.",
            11,
            False,
            TEXT_MUTED,
            0,
        ),
    ]
    _fill_text_frame(body.text_frame, cover_lines, TEXT)

    # ----- One-page agenda (slide 2) -----
    slide_agenda_one_page(prs)

    # ----- Framing -----
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s1, BG_DEEP)
    _gold_rule(s1, Inches(0.42))
    h1 = s1.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12.0), Inches(0.5))
    _fill_text_frame(h1.text_frame, [("HOW TO USE THIS IN THE ROOM", 11, True, GOLD, 0)], TEXT)
    t1 = s1.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.55))
    _fill_text_frame(t1.text_frame, [("Situation  ·  Complication  ·  Implication", 24, True, TEXT, 0)], TEXT)
    _panel(s1, Inches(0.55), Inches(1.55), Inches(12.2), Inches(5.35))
    b1 = s1.shapes.add_textbox(Inches(0.82), Inches(1.75), Inches(11.65), Inches(4.95))
    _fill_text_frame(
        b1.text_frame,
        [
            ("Situation: service is shifting from reactive chat to agentic orchestration — principles must govern autonomy, not only copy.", 12, False, TEXT, 0),
            ("", 6, False, TEXT, 0),
            ("Complication: what helps acquisition (warm tone) can destroy trust in disputes if calibration is wrong.", 12, False, TEXT, 0),
            ("", 6, False, TEXT, 0),
            ("Implication: each principle slide ends with “SCB decisions to lock” — owner, evidence, date.", 12, True, GOLD, 0),
            ("", 8, False, TEXT, 0),
            ("Slide anatomy: gold index bar → principle name as lede → guiding question (purple) → insight line → two panel columns.", 11, False, TEXT_MUTED, 0),
        ],
        TEXT,
    )

    principles = [
        {
            "lede": "Personality & Voice",
            "q": "Who is speaking? And why should I trust them?",
            "insight": "Trust is earned in tone under stress — not in slogans.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Persona = operational spec: diction, pacing, empathy markers; never joke about fraud, death, hardship.", 11, False, TEXT, 1),
                ("• Emotional calibration: after decline or dispute → competence + calm, not cheerleading.", 11, False, TEXT, 1),
                ("• One entity across balance, cards, lending — fragmented “voices” read as untrustworthy bots.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• Card declined → brief apology + factual reason + next step (no celebratory UI).", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• One spine persona + tone-by-state matrix; emoji / humor / Thai–English / voice vs chat rules.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Bank of America — Erica: task-first “virtual financial assistant” framing at US retail scale.", 11, False, TEXT, 1),
                ("• Monzo — plain, human-readable help tone aligned to digital-native brand.", 11, False, TEXT, 1),
                ("• Nubank — informal mass-market voice; contrast case for familiarity vs formality.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• bankofamerica.com/erica", 9, False, TEXT_MUTED, 1),
                ("• monzo.com/help", 9, False, TEXT_MUTED, 1),
            ],
        },
        {
            "lede": "Intent Recognition & Context Awareness",
            "q": "Does it understand me beyond my exact words?",
            "insight": "Understanding must be robust — autonomy multiplies every upstream error.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Intent + entities + confidence; messy phrasing is the default in retail banking.", 11, False, TEXT, 1),
                ("• Session memory contract: coherence vs re-auth for security.", 11, False, TEXT, 1),
                ("• Disambiguation UX: two strong chips beat ten open questions.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• “Weird charge from MERCH*X” → merchant path without exposing cryptic codes.", 11, False, TEXT, 1),
                ("• User corrects amount mid-flow → update slots without re-asking verified answers.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• Journey 2 taxonomy (anomaly, loan hardship, card block) + clarify / act / escalate thresholds.", 11, False, TEXT, 1),
                ("• Cross-session memory + bilingual / code-switch handling.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• BofA Erica — very high volume on bills, transfers, balances → industrial intent coverage.", 11, False, TEXT, 1),
                ("• DBS digibot — authenticated conversational servicing (cards, statements, scams).", 11, False, TEXT, 1),
                ("• Nubank — bounded automation before human (OpenAI case narrative).", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• dbs.com.sg/personal/deposits/bank-with-ease/digibot", 9, False, TEXT_MUTED, 1),
                ("• bankofamerica.com/erica", 9, False, TEXT_MUTED, 1),
                ("• openai.com/index/nubank", 9, False, TEXT_MUTED, 1),
            ],
        },
        {
            "lede": "Conversational Flow Design",
            "q": "Does the conversation move naturally and seamlessly?",
            "insight": "Great service is designed recovery — not a happy path only.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Agentic flows need visible milestones; users accept steps when progress is legible.", 11, False, TEXT, 1),
                ("• Every journey: happy path, branches, cancel/back, timeout, “start over” without penalty.", 11, False, TEXT, 1),
                ("• Friction budget: cap clarification turns before human / callback / secure channel.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• Card block: confirm card → lock → confirmation ID; survive app backgrounding.", 11, False, TEXT, 1),
                ("• Loan hardship: triage → policy hint → handoff packet if human required.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• Canonical diagrams per priority journey + contact-center deflection alignment.", 11, False, TEXT, 1),
                ("• Ownership: in-chat components vs native UI deeplinks.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• DBS digibot — structured servicing + acquisition menus (cards, statements, scams).", 11, False, TEXT, 1),
                ("• BofA Erica — hybrid chat + guided tasks for cards, bills, transfers.", 11, False, TEXT, 1),
                ("• Capital One Eno — proactive decline / recurring-charge style troubleshooting flows.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• dbs.com.sg/personal/deposits/bank-with-ease/digibot", 9, False, TEXT_MUTED, 1),
                ("• capitalone.com/digital/eno", 9, False, TEXT_MUTED, 1),
            ],
        },
        {
            "lede": "Personalization & Relevance",
            "q": "Does it feel like it knows me?",
            "insight": "Personalization without posture is how relevance becomes annoyance.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Shorten time-to-value with verified data; justify proactive content.", 11, False, TEXT, 1),
                ("• Frequency caps; “why am I seeing this?”; servicing-first after failed resolution.", 11, False, TEXT, 1),
                ("• Dynamics vary by segment/channel — not unconstrained model creativity.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• Bill due in 48h → pay-in-thread (high utility) vs loan pitch after chargeback (low).", 11, False, TEXT, 1),
                ("• Subscription price spike → review / cancel merchant (Eno-class hygiene).", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• NBA suppression / cool-down vs journey outcome + sentiment.", 11, False, TEXT, 1),
                ("• Data-use ladder + justification copy + stacked-nudge caps.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• BofA Erica — large proactive “insights” footprint (money moments).", 11, False, TEXT, 1),
                ("• Capital One Eno — recurring / trial / duplicate-charge nudges.", 11, False, TEXT, 1),
                ("• Nubank — high-throughput automated service (verify current guardrails).", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• bankofamerica.com/erica", 9, False, TEXT_MUTED, 1),
                ("• capitalone.com/digital/eno", 9, False, TEXT_MUTED, 1),
            ],
        },
        {
            "lede": "Transparency & Trust Design",
            "q": "Can I see what’s happening?",
            "insight": "In banking, transparency is a product feature — not compliance fine print.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Disclose assistant identity, limits, informational vs actionable answers.", 11, False, TEXT, 1),
                ("• Explainability: merchant, time, FX, fees — plain language before dispute submit.", 11, False, TEXT, 1),
                ("• Risk-tiered confirmations; audit trail for “what did the assistant commit?”", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• Transfer: from/to + amount + fee + timing → explicit confirm.", 11, False, TEXT, 1),
                ("• Agentic plan: “I will check X, then Y” with checkpoints before irreversible calls.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• Disclosure stack + forbidden claims for regulated intents.", 11, False, TEXT, 1),
                ("• Confirmation map + log fields for fraud / dispute.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• DBS digibot — auth before transactional answers (trust gate).", 11, False, TEXT, 1),
                ("• Wise — “show the math” on fees / FX timing (fintech bar).", 11, False, TEXT, 1),
                ("• Monzo — routing to human for urgent / risk cases.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• dbs.com.sg/personal/deposits/bank-with-ease/digibot", 9, False, TEXT_MUTED, 1),
                ("• wise.com/help", 9, False, TEXT_MUTED, 1),
                ("• monzo.com/help", 9, False, TEXT_MUTED, 1),
            ],
        },
        {
            "lede": "Feedback & Learning Loops",
            "q": "Does the system get better over time?",
            "insight": "Learning without governance is silent drift — unacceptable in regulated flows.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Pair explicit (thumbs, CSAT) with implicit (retries, silent exits, escalation codes).", 11, False, TEXT, 1),
                ("• Closed loop: triage → owner → release → post-release monitoring by intent cluster.", 11, False, TEXT, 1),
                ("• Continuous change is gated for regulated intents — approvals + freezes.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• CSAT at resolution milestones only — not every turn.", 11, False, TEXT, 1),
                ("• Weekly top-escalation reasons → prompt / flow change with audit trail.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• Metric suite + RACI for prompt/model updates by journey tier.", 11, False, TEXT, 1),
                ("• One backlog: ACX research + conversational telemetry.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• BofA Erica — multi-billion interactions (public) → industrial learning flywheel.", 11, False, TEXT, 1),
                ("• Nubank — customer automation + agent copilot (dual-loop narrative).", 11, False, TEXT, 1),
                ("• Separate containment from quality — industry norm for mature ops.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• openai.com/index/nubank", 9, False, TEXT_MUTED, 1),
                ("• mckinsey.com (agentic / empathetic CX articles)", 9, False, TEXT_MUTED, 1),
            ],
        },
        {
            "lede": "Human-in-the-Loop & Escalation",
            "q": "Are handoffs to humans seamless and context-complete?",
            "insight": "Escalation quality = context transferred — not containment rate.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Triggers: confidence, sentiment, regulation, fraud, explicit human request.", 11, False, TEXT, 1),
                ("• Warm packet: customer-visible summary + CRM fields agents actually use.", 11, False, TEXT, 1),
                ("• Blended: agent leads; bot supplies checklists only if compliant.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• Fraud: transcript + verified IDs + attempted self-serve → specialist queue.", 11, False, TEXT, 1),
                ("• Callback SMS deep link if voice wait is long.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• Trigger catalog + packet schema + stop-bot-when-human-owns-case.", 11, False, TEXT, 1),
                ("• KPIs must not punish justified handoffs.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Monzo — urgent paths (lost card / fraud) with human access.", 11, False, TEXT, 1),
                ("• Nubank — bounded automation + agent assist (public narrative).", 11, False, TEXT, 1),
                ("• BofA Erica — live help / specialist connection for complex servicing.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• monzo.com/help", 9, False, TEXT_MUTED, 1),
                ("• bankofamerica.com/erica", 9, False, TEXT_MUTED, 1),
            ],
        },
        {
            "lede": "Intuitiveness & Accessibility",
            "q": "Can anyone use this without thinking twice?",
            "insight": "Accessibility is how the product ships — not a separate workstream.",
            "left": [
                ("What this principle means", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• NL-first still uses chips, summaries, structured confirms to cut error.", 11, False, TEXT, 1),
                ("• Multilingual + explicit fallback; plan mixed-language utterances.", 11, False, TEXT, 1),
                ("• Rich media needs text equivalents; voice needs speakable chart summaries.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("Examples", 9, True, GOLD, 0),
                ("• Top intents as quick replies + “None of these”.", 11, False, TEXT, 1),
                ("• Progressive disclosure: one-line answer → expand on tap (reduces text walls).", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("SCB decisions to lock", 9, True, PURPLE, 0),
                ("• Languages, reading level, a11y acceptance criteria for cards/charts.", 11, False, TEXT, 1),
                ("• Voice parity: what cannot be voice-only without confirm.", 11, False, TEXT, 1),
            ],
            "right": [
                ("Exemplars", 9, True, GOLD, 0),
                ("", 3, False, TEXT, 0),
                ("• Category norm: Erica, digibot, Nubank — chat + structured UI hybrid.", 11, False, TEXT, 1),
                ("• Regional leaders — multilingual depth as checklist (e.g. HSBC-style programs).", 11, False, TEXT, 1),
                ("• Fintechs (Revolut / Starling) — speed + clarity patterns; borrow mechanics not voice.", 11, False, TEXT, 1),
                ("", 3, False, TEXT, 0),
                ("References", 9, True, GOLD, 0),
                ("• bankofamerica.com/erica", 9, False, TEXT_MUTED, 1),
                ("• dbs.com.sg/personal/deposits/bank-with-ease/digibot", 9, False, TEXT_MUTED, 1),
            ],
        },
    ]

    # ----- One slide: eight principles + guiding questions (same copy as deep dives) -----
    slide_principles_and_questions(prs, principles)

    for i, p in enumerate(principles, start=1):
        slide_deep_dive(
            prs,
            i,
            p["lede"],
            p["q"],
            p["insight"],
            p["left"],
            p["right"],
        )

    # ----- Sources -----
    sr = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(sr, BG_DEEP)
    _gold_rule(sr, Inches(0.42))
    st = sr.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12.0), Inches(0.5))
    _fill_text_frame(st.text_frame, [("PUBLIC REFERENCES — VERIFY BEFORE EXTERNAL USE", 11, True, GOLD, 0)], TEXT)
    tt = sr.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.45))
    _fill_text_frame(tt.text_frame, [("Selected links", 22, True, TEXT, 0)], TEXT)
    _panel(sr, Inches(0.55), Inches(1.55), Inches(12.2), Inches(5.35))
    sb = sr.shapes.add_textbox(Inches(0.82), Inches(1.75), Inches(11.65), Inches(4.95))
    refs = [
        ("Bank of America — Erica: https://www.bankofamerica.com/erica/", 11, False, TEXT, 0),
        ("DBS — digibot: https://www.dbs.com.sg/personal/deposits/bank-with-ease/digibot", 11, False, TEXT, 0),
        ("Capital One — Eno: https://www.capitalone.com/digital/eno", 11, False, TEXT, 0),
        ("Monzo — Help: https://monzo.com/help", 11, False, TEXT, 0),
        ("Nubank / OpenAI: https://openai.com/index/nubank", 11, False, TEXT, 0),
        ("Wise — Help: https://wise.com/help/", 11, False, TEXT, 0),
        ("McKinsey — empathetic agentic CX: https://www.mckinsey.com/capabilities/operations/our-insights/beyond-the-bot-building-empathetic-customer-experiences-with-agentic-ai", 10, False, TEXT_MUTED, 0),
        ("", 6, False, TEXT, 0),
        ("Design parity: SCB_CXO_Board_Dashboard.html (CXO cockpit). Apply firm template & disclaimers for client distribution.", 10, True, PURPLE, 0),
    ]
    _fill_text_frame(sb.text_frame, refs, TEXT)

    prs.save(OUT)
    print(f"Wrote {OUT}")

    # Standalone one-slide files (same design; skip if open in PowerPoint)
    prs_a = Presentation()
    prs_a.slide_width = Inches(13.333)
    prs_a.slide_height = Inches(7.5)
    slide_agenda_one_page(prs_a)
    try:
        prs_a.save(OUT_AGENDA_ONLY)
        print(f"Wrote {OUT_AGENDA_ONLY}")
    except PermissionError:
        print(f"Skipped (file locked): {OUT_AGENDA_ONLY}")

    prs_q = Presentation()
    prs_q.slide_width = Inches(13.333)
    prs_q.slide_height = Inches(7.5)
    slide_principles_and_questions(prs_q, principles)
    try:
        prs_q.save(OUT_PRINCIPLES_Q_ONLY)
        print(f"Wrote {OUT_PRINCIPLES_Q_ONLY}")
    except PermissionError:
        print(f"Skipped (file locked): {OUT_PRINCIPLES_Q_ONLY}")


if __name__ == "__main__":
    main()
