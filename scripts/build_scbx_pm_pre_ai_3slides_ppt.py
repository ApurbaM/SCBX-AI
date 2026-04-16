from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\Apurba Mukherjee\OneDrive - McKinsey & Company\Desktop\AM\SCBX")
TEMPLATE_PPTX = ROOT / "PM" / "AI Product Management.pptx"
OUT_PPTX = ROOT / "PM" / "SCBx_Product_Management_Guidebook_Pre-AI_3slides.pptx"


@dataclass(frozen=True)
class Stage:
    title: str
    bullets: list[str]
    color: RGBColor


def _remove_all_slides(prs: Presentation) -> None:
    # python-pptx doesn't expose slide deletion publicly; this is the common safe workaround.
    sld_id_lst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    while len(prs.slides) > 0:
        r_id = sld_id_lst[0].rId
        prs.part.drop_rel(r_id)
        del sld_id_lst[0]


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    # Template layouts vary; we use a simple top text box for consistent rendering.
    left = Inches(0.6)
    top = Inches(0.35)
    width = Inches(12.1)
    height = Inches(0.7)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.45))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(16)
        srun.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def _set_bullets(tf, bullets: list[str], font_size_pt: int = 14) -> None:
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size_pt)
        p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)


def build() -> Path:
    # Try to start from template to inherit theme. If OneDrive/permissions prevent reading it,
    # fall back to a clean presentation so we still generate the deck reliably.
    prs: Presentation
    try:
        if not TEMPLATE_PPTX.exists():
            raise FileNotFoundError(f"Template not found: {TEMPLATE_PPTX}")
        prs = Presentation(str(TEMPLATE_PPTX))
        _remove_all_slides(prs)
    except Exception:
        prs = Presentation()

    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Slide 1: Anchor PDLC (Pre-AI)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_title(
        slide1,
        "SCBx Product Delivery Lifecycle (Pre‑AI era)",
        "Anchor view of stages, key activities, and core artefacts (human‑driven delivery).",
    )

    stages = [
        Stage(
            "A. Ideation\n& Viability",
            [
                "Problem statement; customer segments",
                "Market scan / competitive benchmarking",
                "Lean canvas; product vision",
            ],
            RGBColor(0x1F, 0x77, 0xB4),
        ),
        Stage(
            "B. Blueprint\n(Design)",
            [
                "Personas, empathy maps, journey maps",
                "Wireframes; information architecture",
                "Define value prop & OKRs",
            ],
            RGBColor(0x2C, 0xA0, 0x2C),
        ),
        Stage(
            "C. Plan\n& Prioritize",
            [
                "Feature backlog; define MVP",
                "Roadmap (waves/releases)",
                "Prioritization (RICE / MoSCoW)",
            ],
            RGBColor(0xFF, 0x7F, 0x0E),
        ),
        Stage(
            "D. Build\n& Test",
            [
                "Write user stories + acceptance criteria",
                "Sprint goals, cadence, DoR / DoD",
                "Iterative backlog grooming & demos",
            ],
            RGBColor(0x94, 0x67, 0xBD),
        ),
        Stage(
            "E. Launch\n& Scale",
            [
                "GTM plan; launch prep & training",
                "Feedback loops; A/B testing",
                "Operate & support model",
            ],
            RGBColor(0xD6, 0x27, 0x28),
        ),
    ]

    chevron_left = Inches(0.6)
    chevron_top = Inches(1.8)
    chevron_height = Inches(1.1)
    chevron_total_width = Inches(12.1)
    chevron_width = chevron_total_width / len(stages)

    for idx, s in enumerate(stages):
        left = chevron_left + chevron_width * idx
        shape = slide1.shapes.add_shape(MSO_SHAPE.CHEVRON, left, chevron_top, chevron_width, chevron_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = s.color
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = s.title
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # bullets below each chevron
        btop = chevron_top + chevron_height + Inches(0.15)
        bheight = Inches(1.55)
        bbox = slide1.shapes.add_textbox(left, btop, chevron_width, bheight)
        btf = bbox.text_frame
        btf.word_wrap = True
        _set_bullets(btf, s.bullets, font_size_pt=12)

    # Slide 2: PM 101
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide2, "PM 101 (Pre‑AI): What a Product Manager does", "A practical, outcomes‑oriented view for SCBx teams.")

    left_col = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(5.2))
    ltf = left_col.text_frame
    ltf.word_wrap = True
    _set_bullets(
        ltf,
        [
            "Define the product vision aligned to business strategy",
            "Translate user needs into a prioritized roadmap and backlog",
            "Align cross‑functional teams on scope, sequencing, and trade‑offs",
            "Set success metrics (OKRs) and create feedback loops to iterate",
            "Own readiness for launch (GTM) and post‑launch adoption",
        ],
        font_size_pt=16,
    )

    right_col = slide2.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(5.2))
    rtf = right_col.text_frame
    rtf.word_wrap = True
    rtf.clear()
    p0 = rtf.paragraphs[0]
    p0.text = "Core mindsets (pre‑AI)"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    for b in [
        "Customer‑back from personas, empathy & journey maps",
        "Value/feasibility trade‑offs; prioritization discipline",
        "Small increments: MVP → iterate to product‑market fit",
        "Clarity in artefacts: stories + acceptance criteria + DoR/DoD",
        "Cadence and communication: keep stakeholders aligned",
    ]:
        p = rtf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

    # Slide 3: Ways of working + artefacts
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide3, "Ways of working (Pre‑AI): ceremonies, roles, and artefacts", "The minimum operating rhythm to deliver consistently.")

    # Ceremonies
    cer = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(4.0), Inches(5.4))
    cer.fill.solid()
    cer.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    cer.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    ctf = cer.text_frame
    ctf.clear()
    ct = ctf.paragraphs[0]
    ct.text = "Ceremonies / cadence"
    ct.font.bold = True
    ct.font.size = Pt(16)
    for b in [
        "Sprint planning (commit scope)",
        "Backlog grooming/refinement",
        "Daily stand‑up (remove blockers)",
        "Sprint review/demo (share progress)",
        "Retro (improve ways of working)",
        "User testing playbacks (learn fast)",
    ]:
        p = ctf.add_paragraph()
        p.text = b
        p.font.size = Pt(13)

    # Roles
    roles = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.05), Inches(1.75), Inches(3.75), Inches(5.4))
    roles.fill.solid()
    roles.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    roles.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    rtf2 = roles.text_frame
    rtf2.clear()
    rt = rtf2.paragraphs[0]
    rt.text = "Core roles"
    rt.font.bold = True
    rt.font.size = Pt(16)
    for b in [
        "PM / Product Owner",
        "Tech lead / Architect",
        "Engineers (FE/BE)",
        "Product designer (UX/UI)",
        "QA / testing (as needed)",
        "Scrum master (as needed)",
        "Stakeholders / SMEs",
    ]:
        p = rtf2.add_paragraph()
        p.text = b
        p.font.size = Pt(13)

    # Artefacts
    art = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(1.75), Inches(3.55), Inches(5.4))
    art.fill.solid()
    art.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    art.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    atf = art.text_frame
    atf.clear()
    at = atf.paragraphs[0]
    at.text = "Core artefacts"
    at.font.bold = True
    at.font.size = Pt(16)
    for b in [
        "Lean canvas + product vision",
        "Personas / empathy maps / journey maps",
        "Value proposition + OKRs",
        "Roadmap (waves/releases)",
        "Backlog + user stories (INVEST) + acceptance criteria",
        "Definition of Ready / Done",
        "Release plan + GTM checklist",
    ]:
        p = atf.add_paragraph()
        p.text = b
        p.font.size = Pt(13)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    return OUT_PPTX


if __name__ == "__main__":
    out = build()
    print(str(out))

